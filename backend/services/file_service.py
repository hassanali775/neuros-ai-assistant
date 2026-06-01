import uuid
import os
import mimetypes
from pathlib import Path
from fastapi import UploadFile, HTTPException
from sqlalchemy.ext.asyncio import AsyncSession
from sqlalchemy import select
from models.database import FileAttachment
from models.schemas import FileUploadResponse
from config import get_settings

# Import specialized file parsers
from pypdf import PdfReader
try:
    from pptx import Presentation
except ImportError:
    Presentation = None

settings = get_settings()

ALLOWED_MIME_TYPES = {
    "text/plain",
    "text/markdown",
    "text/csv",
    "application/pdf",
    "application/json",
    "application/javascript",
    "text/html",
    "text/css",
    "text/x-python",
    "application/x-python-code",
    "image/jpeg",
    "image/png",
    "image/gif",
    "image/webp",
    "application/vnd.ms-powerpoint",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation"
}


class FileService:
    def __init__(self):
        self.upload_dir = Path(settings.upload_dir)
        self.upload_dir.mkdir(parents=True, exist_ok=True)
        self.max_size = settings.max_upload_size_mb * 1024 * 1024

    async def upload(self, db: AsyncSession, file: UploadFile) -> FileUploadResponse:
        content = await file.read()

        if len(content) > self.max_size:
            raise HTTPException(
                status_code=413,
                detail=f"File too large. Maximum size: {settings.max_upload_size_mb}MB"
            )

        mime_type = file.content_type or "application/octet-stream"
        if mime_type == "application/octet-stream":
            guessed, _ = mimetypes.guess_type(file.filename or "")
            mime_type = guessed or "application/octet-stream"

        if mime_type not in ALLOWED_MIME_TYPES:
            raise HTTPException(
                status_code=415,
                detail=f"Unsupported file type: {mime_type}"
            )

        file_id = str(uuid.uuid4())
        ext = Path(file.filename or "file").suffix
        stored_name = f"{file_id}{ext}"
        file_path = self.upload_dir / stored_name

        with open(file_path, "wb") as f:
            f.write(content)

        attachment = FileAttachment(
            id=file_id,
            message_id="pending",
            filename=stored_name,
            original_name=file.filename or "unnamed",
            mime_type=mime_type,
            size_bytes=len(content),
        )
        db.add(attachment)
        await db.commit()

        return FileUploadResponse(
            id=file_id,
            filename=stored_name,
            original_name=file.filename or "unnamed",
            mime_type=mime_type,
            size_bytes=len(content),
        )

    async def read_text(self, file_id: str) -> str | None:
        """Safely extracts clear text layers from text files, PDFs, and PPTX slideshows."""
        target_file = None
        for f in self.upload_dir.iterdir():
            if f.stem == file_id:
                target_file = f
                break

        if not target_file or not target_file.exists():
            return None

        file_ext = target_file.suffix.lower()

        try:
            # 1. HANDLE BINARY PDFs
            if file_ext == ".pdf":
                reader = PdfReader(target_file)
                extracted_text = []
                for page in reader.pages:
                    text_layer = page.extract_text()
                    if text_layer:
                        extracted_text.append(text_layer)
                return "\n".join(extracted_text).strip() if extracted_text else None

            # 2. HANDLE BINARY POWERPOINT SLIDES
            if file_ext in [".pptx", ".ppt"]:
                if not Presentation:
                    print("──> [FILE ENGINE] python-pptx library is missing.")
                    return None
                prs = Presentation(target_file)
                extracted_text = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            extracted_text.append(shape.text.strip())
                return "\n".join(extracted_text).strip() if extracted_text else None

            # 3. HANDLE STANDARD TEXT/CODE FORMATS
            return target_file.read_text(encoding="utf-8", errors="replace")

        except Exception as e:
            print(f"──> [FILE ENGINE ERROR] Critical parsing block drop for {target_file.name}: {e}")
            return None

    async def get_by_id(self, db: AsyncSession, file_id: str) -> FileAttachment | None:
        stmt = select(FileAttachment).where(FileAttachment.id == file_id)
        result = await db.execute(stmt)
        return result.scalar_one_or_none()

    async def link_to_message(self, db: AsyncSession, file_ids: list[str], message_id: str) -> None:
        for file_id in file_ids:
            stmt = select(FileAttachment).where(FileAttachment.id == file_id)
            result = await db.execute(stmt)
            attachment = result.scalar_one_or_none()
            if attachment:
                attachment.message_id = message_id
        await db.commit()

    def get_file_path(self, filename: str) -> Path:
        return self.upload_dir / filename


def get_file_service() -> FileService:
    return FileService()